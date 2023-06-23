import os
from os import path
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Cm

from utils import ASSETS_DIR


def add_image_slide(prs, img_path):
    # Add a slide with your specific layout
    slide_layout = prs.slide_layouts[19]
    slide = prs.slides.add_slide(slide_layout)

    # Get the title placeholder
    title_placeholder = slide.shapes.title

    # Set the title to the image name (without the extension)
    image_name = path.splitext(path.basename(img_path))[
        0
    ]  # This gets the file name without the extension

    # Format the image name: replace underscores with spaces and capitalize each word
    title_text = " ".join(word.capitalize() for word in image_name.split("_"))
    title_placeholder.text = title_text

    # Load image
    with open(img_path, "rb") as image_stream:
        image = slide.shapes.add_picture(image_stream, 0, 0)

    # Set a maximum height for the image and maintain the aspect ratio
    if image.height > Cm(13):
        aspect_ratio = image.width / image.height
        image.height = Cm(13)
        image.width = int(image.height * aspect_ratio)

    # Calculate the picture's centered position
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    left = int((slide_width - image.width) / 2)
    top = int((slide_height - image.height) / 2)

    # Position the picture in the center
    image.left = left
    image.top = top + Inches(0.5)


def generate_ppt_from_plots(results_dir: Path):
    pptx_path = results_dir / "snowlit_plots.pptx"
    if pptx_path.exists():
        print("Plots are already exported at", pptx_path, "skipping...")
        return

    # Create presentation
    prs = Presentation(ASSETS_DIR / "template.pptx")
    plots_dir = results_dir / "plots"
    # Get a sorted list of all .png files in the directory, sorted by creation time
    files = sorted(
        (
            results_dir / "plots" / filename
            for filename in os.listdir(plots_dir)
            if filename.endswith(".png")
        ),
        key=os.path.getctime,
    )

    # Add a slide with your specific layout
    slide_layout = prs.slide_layouts[
        5
    ]  # Assume a layout that contains a text box or placeholder
    slide = prs.slides.add_slide(slide_layout)

    # Get the title placeholder
    title_placeholder = slide.shapes.title

    # Name the title Search Query
    title_text = "Search Query"
    title_placeholder.text = title_text

    # Access the text placeholder to enter search querry
    text_box = slide.shapes.placeholders[1]  # Index may vary depending on layout
    text_frame = text_box.text_frame

    # Clear any existing text in the placeholder
    text_frame.clear()
    text_content = "LCA search bla bla"

    # Add and format the new text
    p = text_frame.add_paragraph()
    p.text = text_content

    # Loop over each file in the sorted list
    for full_path in files:
        add_image_slide(prs, full_path)  # add slide with the image

    # Save the presentation

    prs.save(pptx_path)

    print("All plots have been exported at ", pptx_path, " successfully!")
